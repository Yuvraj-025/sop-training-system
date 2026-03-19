import os
import uuid
import json
import logging
from pathlib import Path
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

OUTPUT_DIR = Path("/app/output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

app = FastAPI(
    title="SOP Presentation Service",
    description="Generates downloadable PowerPoint presentations from AI-processed SOP data",
    version="1.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Color Palette ──────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
ACCENT    = RGBColor(0x6C, 0x63, 0xFF)   # Purple accent
ACCENT2   = RGBColor(0x00, 0xD4, 0xFF)   # Cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xB0, 0xB8, 0xD4)
CARD_BG   = RGBColor(0x1A, 0x23, 0x3E)


class AIOutput(BaseModel):
    summary: dict
    training: dict
    quiz: dict
    filename: str = "SOP_Training"


def set_slide_background(slide, prs, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_accent_bar(slide, top=1.2):
    """Add a horizontal accent line."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(top), Inches(12.33), Inches(0.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def build_title_slide(prs: Presentation, summary: dict):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    # Top accent strip
    strip = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.8))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    # Title
    add_textbox(slide, summary.get("title", "SOP Training Program"),
                0.5, 1.5, 12.33, 1.5, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, summary.get("overview", "")[:200],
                1.0, 3.2, 11.33, 1.2, font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Department pill
    dept = summary.get("department", "")
    if dept:
        add_textbox(slide, f"📂  {dept}",
                    4.5, 4.6, 4.33, 0.5, font_size=14, color=ACCENT2, align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, "AI-Generated Training Package  •  Confidential",
                0.5, 6.8, 12.33, 0.4, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def build_summary_slide(prs, summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs, DARK_BG)

    add_textbox(slide, "📋  SOP Overview", 0.5, 0.3, 12, 0.7, font_size=26, bold=True, color=WHITE)
    add_accent_bar(slide, 1.15)

    # Purpose + Scope
    add_textbox(slide, "PURPOSE", 0.5, 1.35, 5.5, 0.4, font_size=11, bold=True, color=ACCENT2)
    add_textbox(slide, summary.get("purpose", ""), 0.5, 1.7, 5.8, 1.0, font_size=13, color=LIGHT_GRAY)

    add_textbox(slide, "SCOPE", 7.0, 1.35, 5.5, 0.4, font_size=11, bold=True, color=ACCENT2)
    add_textbox(slide, summary.get("scope", ""), 7.0, 1.7, 5.8, 1.0, font_size=13, color=LIGHT_GRAY)

    # Key points
    add_textbox(slide, "KEY POINTS", 0.5, 2.9, 12, 0.4, font_size=11, bold=True, color=ACCENT2)
    key_points = summary.get("key_points", [])[:6]
    for i, pt in enumerate(key_points):
        col = 0.5 if i % 2 == 0 else 6.7
        row = 3.35 + (i // 2) * 0.75
        add_textbox(slide, f"✦  {pt}", col, row, 5.8, 0.65, font_size=12, color=WHITE)


def build_sections_slide(prs, summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs, DARK_BG)

    add_textbox(slide, "📄  Document Sections", 0.5, 0.3, 12, 0.7, font_size=26, bold=True, color=WHITE)
    add_accent_bar(slide)

    sections = summary.get("sections", [])[:6]
    for i, sec in enumerate(sections):
        col = 0.5 if i % 2 == 0 else 6.7
        row = 1.4 + (i // 2) * 1.5

        # Card background
        card = slide.shapes.add_shape(1, Inches(col - 0.1), Inches(row - 0.1), Inches(5.8), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT

        add_textbox(slide, sec.get("heading", ""), col, row, 5.5, 0.4, font_size=13, bold=True, color=ACCENT2)
        add_textbox(slide, sec.get("summary", "")[:130], col, row + 0.45, 5.5, 0.75, font_size=11, color=LIGHT_GRAY)


def build_training_module_slide(prs, module: dict, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs, DARK_BG)

    add_textbox(slide, f"🎓  Module {idx}/{total}: {module.get('title', '')}", 0.5, 0.3, 12.5, 0.7,
                font_size=22, bold=True, color=WHITE)
    add_accent_bar(slide)

    add_textbox(slide, f"Objective: {module.get('objective', '')}", 0.5, 1.25, 12.33, 0.55,
                font_size=13, color=ACCENT2)

    steps = module.get("steps", [])[:5]
    add_textbox(slide, "STEPS", 0.5, 1.9, 7.5, 0.35, font_size=11, bold=True, color=LIGHT_GRAY)
    for i, step in enumerate(steps):
        action = step.get("action", "")
        details = step.get("details", "")
        add_textbox(slide, f"{step.get('step_number', i+1)}.  {action}", 0.5, 2.3 + i * 0.75, 7.5, 0.4,
                    font_size=12, bold=True, color=WHITE)
        add_textbox(slide, details[:110], 1.0, 2.65 + i * 0.75, 7.0, 0.35, font_size=11, color=LIGHT_GRAY)

    # Tips panel
    tips = module.get("tips", [])[:3]
    if tips:
        add_textbox(slide, "💡 TIPS", 9.0, 1.9, 3.8, 0.35, font_size=11, bold=True, color=ACCENT2)
        for i, tip in enumerate(tips):
            add_textbox(slide, f"• {tip[:80]}", 9.0, 2.3 + i * 0.8, 3.8, 0.7, font_size=11, color=LIGHT_GRAY)


def build_quiz_slide(prs, question: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs, DARK_BG)

    qnum = question.get("question_number", "?")
    diff = question.get("difficulty", "medium")
    diff_color = RGBColor(0x00, 0xD4, 0x8A) if diff == "easy" else (ACCENT if diff == "medium" else RGBColor(0xFF, 0x6B, 0x6B))

    add_textbox(slide, f"❓  Quiz — Question {qnum}", 0.5, 0.3, 10, 0.6, font_size=24, bold=True, color=WHITE)
    add_textbox(slide, f"  {diff.upper()}  ", 11.0, 0.35, 2.0, 0.5, font_size=12, bold=True, color=diff_color, align=PP_ALIGN.CENTER)
    add_accent_bar(slide)

    add_textbox(slide, question.get("question", ""), 0.5, 1.3, 12.33, 0.9, font_size=15, bold=True, color=WHITE)

    options = question.get("options", {})
    colors = {"A": RGBColor(0x6C, 0x63, 0xFF), "B": RGBColor(0x00, 0xD4, 0xFF),
              "C": RGBColor(0xFF, 0xA5, 0x00), "D": RGBColor(0xFF, 0x6B, 0x6B)}
    top = 2.4
    for key, text in options.items():
        card = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(12.33), Inches(0.65))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = colors.get(key, ACCENT)
        add_textbox(slide, f"  {key}.  {text}", 0.6, top + 0.1, 12.0, 0.5, font_size=13, color=WHITE)
        top += 0.8

    answer = question.get("correct_answer", "")
    explanation = question.get("explanation", "")
    add_textbox(slide, f"✅  Answer: {answer}  —  {explanation[:160]}", 0.5, 6.2, 12.33, 0.6,
                font_size=12, color=RGBColor(0x00, 0xD4, 0x8A))


def build_thank_you_slide(prs, summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs, DARK_BG)

    strip = slide.shapes.add_shape(1, Inches(0), Inches(6.6), Inches(13.33), Inches(0.9))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    add_textbox(slide, "🎉", 5.5, 1.5, 2.33, 0.8, font_size=40, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Training Complete!", 0.5, 2.5, 12.33, 0.9, font_size=36, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "You are now ready to apply this SOP in your daily work.",
                1.0, 3.5, 11.33, 0.7, font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    notes = summary.get("compliance_notes", "")
    if notes:
        add_textbox(slide, f"📌  {notes[:200]}", 1.5, 4.4, 10.33, 0.8, font_size=12, color=ACCENT2, align=PP_ALIGN.CENTER)


def generate_pptx(data: AIOutput) -> Path:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    summary  = data.summary
    training = data.training
    quiz     = data.quiz

    build_title_slide(prs, summary)
    build_summary_slide(prs, summary)
    build_sections_slide(prs, summary)

    modules = training.get("modules", [])
    for i, mod in enumerate(modules, 1):
        build_training_module_slide(prs, mod, i, len(modules))

    questions = quiz.get("questions", [])
    for q in questions:
        build_quiz_slide(prs, q)

    build_thank_you_slide(prs, summary)

    job_id = str(uuid.uuid4())
    output_path = OUTPUT_DIR / f"{job_id}.pptx"
    prs.save(str(output_path))
    logger.info(f"Saved presentation: {output_path}")
    return output_path, job_id


@app.get("/health")
async def health():
    return {"status": "healthy", "service": "presentation"}


@app.post("/generate")
async def generate(data: AIOutput):
    """Generate a PowerPoint presentation from AI output."""
    try:
        output_path, job_id = generate_pptx(data)
        return {
            "job_id": job_id,
            "filename": f"{data.filename.replace(' ', '_')}_Training.pptx",
            "download_url": f"/download/{job_id}",
            "slide_count": len(data.training.get("modules", [])) + len(data.quiz.get("questions", [])) + 4,
        }
    except Exception as e:
        logger.error(f"Presentation generation failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/download/{job_id}")
async def download(job_id: str):
    """Download the generated .pptx file."""
    file_path = OUTPUT_DIR / f"{job_id}.pptx"
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Presentation not found.")
    return FileResponse(
        path=str(file_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{job_id}_training.pptx",
    )
